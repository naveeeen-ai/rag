import os
from pathlib import Path
from typing import List
import shutil
from io import BytesIO

from langchain_chroma import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

# PDF OCR
import fitz  # PyMuPDF
from PIL import Image
import pytesseract

# PPTX
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from get_embedding_function import get_embedding_function
from dotenv import load_dotenv
load_dotenv()


CHROMA_DIR = os.getenv("CHROMA_DIR", "chroma_db")
DATA_DIR = os.getenv("DATA_DIR", "data")


def is_tesseract_available() -> bool:
    try:
        _ = pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def list_input_files(directory: str) -> List[Path]:
    base = Path(directory)
    files: List[Path] = []
    for pattern in ("**/*.pdf", "**/*.pptx", "**/*.ppt"):
        files.extend([p for p in base.glob(pattern) if p.is_file()])
    # Warn on legacy .ppt
    ppt_files = [p for p in files if p.suffix.lower() == ".ppt"]
    if ppt_files:
        print("Warning: .ppt detected. Please convert to .pptx or .pdf for best results.")
    # Only process pdf and pptx
    return sorted([p for p in files if p.suffix.lower() in {".pdf", ".pptx"}])


def _ocr_image_bytes(image_bytes: bytes) -> str:
    try:
        img = Image.open(BytesIO(image_bytes)).convert("RGB")
        return pytesseract.image_to_string(img)
    except Exception:
        return ""


def _ocr_pixmap(page, zoom: float = 2.0) -> str:
    try:
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        return pytesseract.image_to_string(img)
    except Exception:
        return ""


def load_pdf_documents(pdf_path: Path, ocr_enabled: bool) -> List[Document]:
    docs: List[Document] = []
    with fitz.open(str(pdf_path)) as doc:
        for i, page in enumerate(doc):
            text = page.get_text("text") or ""
            ocr_text = ""
            if ocr_enabled and (not text.strip()):
                ocr_text = _ocr_pixmap(page)
            content = (text + ("\n\n" + ocr_text if ocr_text.strip() else "")).strip()
            if not content:
                continue
            docs.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": str(pdf_path),
                        "page": i + 1,
                        "type": "pdf",
                    },
                )
            )
    return docs


def load_pptx_documents(pptx_path: Path, ocr_enabled: bool) -> List[Document]:
    prs = Presentation(str(pptx_path))
    docs: List[Document] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        # Text from shapes
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                tb = shape.text_frame
                if tb is not None:
                    texts.append(tb.text or "")
        # OCR on images
        if ocr_enabled:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_bytes = shape.image.blob
                        ocr_text = _ocr_image_bytes(image_bytes)
                        if ocr_text.strip():
                            texts.append(ocr_text)
                    except Exception:
                        continue
        content = "\n".join([t for t in texts if t and t.strip()]).strip()
        if not content:
            continue
        docs.append(
            Document(
                page_content=content,
                metadata={
                    "source": str(pptx_path),
                    "page": idx,
                    "type": "pptx",
                },
            )
        )
    return docs


def split_documents(documents: List[Document]) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=150,
        length_function=len,
        is_separator_regex=False,
    )
    return splitter.split_documents(documents)


def reset_vector_db():
    if os.path.isdir(CHROMA_DIR):
        print(f"Removing existing vector DB at '{CHROMA_DIR}'...")
        shutil.rmtree(CHROMA_DIR, ignore_errors=True)


def main():
    reset_vector_db()
    os.makedirs(CHROMA_DIR, exist_ok=True)

    input_files = list_input_files(DATA_DIR)
    if not input_files:
        print(f"No PDFs/PPTX found in '{DATA_DIR}'. Add files and run again.")
        return

    print(f"Found {len(input_files)} file(s). Loading...")
    use_ocr = is_tesseract_available()
    if not use_ocr:
        print("Note: Tesseract not found. OCR will be skipped. Install Tesseract to extract text from images.")

    documents: List[Document] = []
    for path in input_files:
        if path.suffix.lower() == ".pdf":
            documents.extend(load_pdf_documents(path, use_ocr))
        elif path.suffix.lower() == ".pptx":
            documents.extend(load_pptx_documents(path, use_ocr))

    if not documents:
        print("No textual content extracted. Nothing to index.")
        return

    print(f"Loaded {len(documents)} documents. Splitting...")
    chunks = split_documents(documents)
    print(f"Created {len(chunks)} chunks. Embedding and adding to Chroma...")

    embeddings = get_embedding_function()

    vectordb = Chroma(
        collection_name=os.getenv("COLLECTION_NAME", "rag_collection"),
        persist_directory=CHROMA_DIR,
        embedding_function=embeddings,
    )

    vectordb.add_documents(chunks)
    print(f"Ingestion complete. DB persisted at '{CHROMA_DIR}'.")


if __name__ == "__main__":
    main()
